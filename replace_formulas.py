from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE

prs = Presentation(r'D:\MCP\MCP_Academic_Ultimate.pptx')

# Mapping: slide number (1-based) -> list of formula image paths to insert
formula_map = {
    6: [r'D:\MCP\assets\formulas\eq_latency.png'],
    7: [r'D:\MCP\assets\formulas\eq_prob_chain.png'],
    8: [r'D:\MCP\assets\formulas\eq_success_bound.png'],
    9: [r'D:\MCP\assets\formulas\eq_cost.png'],
    10: [r'D:\MCP\assets\formulas\eq_throughput.png'],
    11: [r'D:\MCP\assets\formulas\eq_tradeoff.png'],
}

for slide_num, new_imgs in formula_map.items():
    slide = prs.slides[slide_num - 1]
    
    # Remove ALL existing pictures from this slide (old formula + chart images)
    pics_to_remove = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            pics_to_remove.append(shape)
    
    for pic in pics_to_remove:
        sp = pic._element
        sp.getparent().remove(sp)
        print(f'  Removed old image from slide {slide_num}')
    
    # Insert new formula image
    for img_path in new_imgs:
        slide.shapes.add_picture(img_path, Inches(0.9), Inches(3.1), Inches(8.0))
        print(f'  Inserted {img_path.split(chr(92))[-1]} -> slide {slide_num}')

# Re-insert chart images on slides that also need them
chart_inserts = {
    6: r'D:\MCP\assets\charts\latency_vs_tools.png',
    8: r'D:\MCP\assets\charts\success_rate_bar.png',
    9: r'D:\MCP\assets\charts\cost_vs_tokens.png',
    11: r'D:\MCP\assets\charts\tradeoff_scatter.png',
}

for slide_num, chart_path in chart_inserts.items():
    slide = prs.slides[slide_num - 1]
    slide.shapes.add_picture(chart_path, Inches(2.2), Inches(5.0), Inches(5.5))
    print(f'  Inserted chart {chart_path.split(chr(92))[-1]} -> slide {slide_num}')

prs.save(r'D:\MCP\MCP_Academic_Ultimate.pptx')
print('SAVE OK - all formula images replaced with LaTeX-rendered versions')
