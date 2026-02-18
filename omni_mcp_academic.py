#!/usr/bin/env python3
"""omni_mcp v2 (academic-style mirror)
Unified MCP server implementation for multimodal engineering workflows.
The service integrates office-document automation, raster/vector graphics,
media transcoding, and 3D/CAD tool orchestration under a single tool API.
This file is a style-refined copy of ``d:\\omni_mcp.py`` with behavior preserved.
"""
import glob
import json
import os
import re
import shutil
import subprocess
import tempfile
from pathlib import Path
from typing import Optional, Sequence, Tuple
from mcp.server.fastmcp import FastMCP
mcp = FastMCP("omni_mcp")
# ========== CONFIG ==========
BLENDER = r"D:\Blender\blender.exe"
WD = Path(tempfile.gettempdir()) / "omni_mcp"
WD.mkdir(exist_ok=True)
CF = subprocess.CREATE_NO_WINDOW if os.name == "nt" else 0
def _find(*candidates: str) -> Optional[str]:
    """Return the first executable found from glob candidates."""
    for candidate in candidates:
        for resolved in glob.glob(candidate):
            if os.path.isfile(resolved):
                return resolved
    return None
MATLAB = _find(r"C:\Program Files\MATLAB\*\bin\matlab.exe", r"D:\MATLAB\*\bin\matlab.exe") or "matlab"
FFMPEG = _find(r"C:\Users\*\AppData\Local\Microsoft\WinGet\Packages\*ffmpeg*\ffmpeg*.exe") or "ffmpeg"
GIMP = (
    _find(r"D:\GIMP*\bin\gimp-console-*.exe", r"C:\Program Files\GIMP*\bin\gimp-console-*.exe")
    or _find(r"D:\GIMP*\bin\gimp-*.exe", r"C:\Program Files\GIMP*\bin\gimp-*.exe")
    or "gimp"
)
INKSCAPE = _find(r"D:\Inkscape*\bin\inkscape.exe", r"C:\Program Files\Inkscape*\bin\inkscape.exe") or "inkscape"
FREECAD = (
    _find(
        r"D:\FreeCAD*\bin\FreeCADCmd.exe",
        r"C:\Program Files\FreeCAD*\bin\FreeCADCmd.exe",
        r"D:\FreeCAD*\bin\freecadcmd.exe",
        r"C:\Program Files\FreeCAD*\bin\freecadcmd.exe",
    )
    or "FreeCADCmd"
)
GODOT = _find(r"D:\Godot*\Godot*.exe", r"C:\Godot*\Godot*.exe") or "godot"
def R(path: str) -> Path:
    """Resolve a relative path under the MCP working directory."""
    p = Path(path)
    return p if p.is_absolute() else WD / p
def J(ok: bool = True, **payload) -> str:
    """Serialize MCP responses with UTF-8 safe settings."""
    return json.dumps({"ok": ok, **payload}, ensure_ascii=False, default=str)
def _run(
    cmd: Sequence[str] | str,
    timeout: int = 120,
    shell: bool = False,
    cwd: Optional[str] = None,
) -> Tuple[str, str, int]:
    """Run a subprocess and clamp output length for MCP transport stability."""
    result = subprocess.run(
        cmd,
        capture_output=True,
        text=True,
        timeout=timeout,
        shell=shell,
        cwd=cwd or str(WD),
        creationflags=CF,
        errors="replace",
        stdin=subprocess.DEVNULL,
    )
    stdout = result.stdout or ""
    stderr = result.stderr or ""
    if len(stdout) > 4000:
        stdout = stdout[:2000] + "\n...(truncated)...\n" + stdout[-2000:]
    if len(stderr) > 2000:
        stderr = stderr[-2000:]
    return stdout, stderr, result.returncode
# --- PPTX SUBSYSTEM: Presentation I/O and editing ---
@mcp.tool(name="pptx_create")
async def pptx_create(path:str,slides:str="[]",template:str="")->str:
    """创建PPT
    slides: JSON数组, 每项:
      {"title":"","content":"","layout":1,"notes":"",
       "font_size":18,"font_color":"000000","bold":false,
       "bg_color":"FFFFFF","images":[{"path":"x.png","left":1,"top":1,"width":5}]}
    layout: 0=标题页 1=标题+内容 2=节标题 5=空白 6=仅内容
    template: 可选,模板pptx路径"""
    from pptx import Presentation
    from pptx.util import Inches,Pt,Emu
    from pptx.dml.color import RGBColor
    try:
        prs=Presentation(str(R(template))) if template else Presentation()
        for s in json.loads(slides):
            ly=prs.slide_layouts[s.get("layout",1)]
            sl=prs.slides.add_slide(ly)
            ph=list(sl.placeholders)
            # 背景色
            if s.get("bg_color"):
                bg=sl.background
                fill=bg.fill
                fill.solid()
                fill.fore_color.rgb=RGBColor.from_string(s["bg_color"])
            # 标题
            if ph and s.get("title"):
                ph[0].text=s["title"]
                if s.get("font_size"):
                    for run in ph[0].text_frame.paragraphs[0].runs:
                        run.font.size=Pt(s["font_size"])
                if s.get("font_color"):
                    for run in ph[0].text_frame.paragraphs[0].runs:
                        run.font.color.rgb=RGBColor.from_string(s["font_color"])
                if s.get("bold"):
                    for run in ph[0].text_frame.paragraphs[0].runs:
                        run.font.bold=True
            # 内容
            if len(ph)>1 and s.get("content"):
                ph[1].text=s["content"]
            # 备注
            if s.get("notes"):
                sl.notes_slide.notes_text_frame.text=s["notes"]
            # 图片
            for img in s.get("images",[]):
                sl.shapes.add_picture(str(R(img["path"])),
                    Inches(img.get("left",1)),Inches(img.get("top",1)),
                    Inches(img.get("width",5)))
        p=str(R(path))
        prs.save(p)
        return J(path=p,count=len(prs.slides))
    except Exception as e:
        return J(False,err=str(e))
@mcp.tool(name="pptx_read")
async def pptx_read(path:str)->str:
    """读取PPT全部文本"""
    from pptx import Presentation
    try:
        prs=Presentation(str(R(path)))
        out=[]
        for i,sl in enumerate(prs.slides):
            texts=[sh.text for sh in sl.shapes if sh.has_text_frame]
            notes=""
            try:
                notes=sl.notes_slide.notes_text_frame.text
            except:
                pass
            out.append({"slide":i+1,"texts":texts,"notes":notes})
        return J(slides=out,count=len(out))
    except Exception as e:
        return J(False,err=str(e))
@mcp.tool(name="pptx_edit")
async def pptx_edit(path:str,ops:str="[]")->str:
    """编辑PPT
    ops: JSON数组:
      {"slide":1,"placeholder":0,"text":"新文本"}
      {"slide":1,"placeholder":0,"font_size":24,"bold":true,"font_color":"FF0000"}
      {"slide":1,"add_image":"图片路径","left":1,"top":1,"width":5,"height":3}
      {"slide":1,"add_textbox":"文本","left":1,"top":1,"width":4,"height":1,"font_size":14}
      {"slide":1,"add_shape":"rect","left":1,"top":1,"width":3,"height":2,"fill":"0070C0"}
      {"slide":1,"delete":true}
    slide从1开始"""
    from pptx import Presentation
    from pptx.util import Inches,Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    try:
        p=str(R(path))
        prs=Presentation(p)
        for op in json.loads(ops):
            si=op["slide"]-1
            sl=prs.slides[si]
            if op.get("delete"):
                rId=prs.slides._sldIdLst[si].rId
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[si]
            elif "text" in op:
                ph=sl.placeholders[op.get("placeholder",0)]
                ph.text=op["text"]
                if op.get("font_size") or op.get("bold") or op.get("font_color"):
                    for para in ph.text_frame.paragraphs:
                        for run in para.runs:
                            if op.get("font_size"):
                                run.font.size=Pt(op["font_size"])
                            if op.get("bold"):
                                run.font.bold=True
                            if op.get("font_color"):
                                run.font.color.rgb=RGBColor.from_string(op["font_color"])
            elif "add_image" in op:
                kw={"width":Inches(op.get("width",5))}
                if op.get("height"):
                    kw["height"]=Inches(op["height"])
                sl.shapes.add_picture(str(R(op["add_image"])),
                    Inches(op.get("left",1)),Inches(op.get("top",1)),**kw)
            elif "add_textbox" in op:
                from pptx.util import Inches as In
                txBox=sl.shapes.add_textbox(In(op.get("left",1)),In(op.get("top",1)),
                    In(op.get("width",4)),In(op.get("height",1)))
                tf=txBox.text_frame
                tf.text=op["add_textbox"]
                if op.get("font_size"):
                    for r in tf.paragraphs[0].runs:
                        r.font.size=Pt(op["font_size"])
            elif "add_shape" in op:
                shapes_map={"rect":MSO_SHAPE.RECTANGLE,"oval":MSO_SHAPE.OVAL,
                           "triangle":MSO_SHAPE.ISOSCELES_TRIANGLE,"arrow":MSO_SHAPE.RIGHT_ARROW,
                           "star":MSO_SHAPE.STAR_5_POINT,"diamond":MSO_SHAPE.DIAMOND}
                st=shapes_map.get(op["add_shape"],MSO_SHAPE.RECTANGLE)
                sh=sl.shapes.add_shape(st,Inches(op.get("left",1)),Inches(op.get("top",1)),
                    Inches(op.get("width",3)),Inches(op.get("height",2)))
                if op.get("fill"):
                    sh.fill.solid()
                    sh.fill.fore_color.rgb=RGBColor.from_string(op["fill"])
        prs.save(p)
        return J(path=p)
    except Exception as e:
        return J(False,err=str(e))
# --- DOCX SUBSYSTEM: Document generation and replacement ---
@mcp.tool(name="docx_create")
async def docx_create(path:str,content:str="[]",template:str="")->str:
    """创建Word文档
    content: JSON数组:
      {"type":"heading","text":"","level":1}
      {"type":"para","text":"","bold":false,"italic":false,"size":12,"color":"000000","align":"left"}
      {"type":"table","rows":[["a","b"],["c","d"]],"header":true,"col_widths":[3,5]}
      {"type":"image","path":"","width":5}
      {"type":"list","items":["a","b","c"],"style":"bullet"|"number"}
      {"type":"toc"} (目录占位)
      {"type":"hr"} (水平线)
      {"type":"break"}
    template: 可选模板docx路径"""
    from docx import Document
    from docx.shared import Inches,Pt,RGBColor,Cm
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    try:
        doc=Document(str(R(template))) if template else Document()
        aligns={"left":WD_ALIGN_PARAGRAPH.LEFT,"center":WD_ALIGN_PARAGRAPH.CENTER,
                "right":WD_ALIGN_PARAGRAPH.RIGHT,"justify":WD_ALIGN_PARAGRAPH.JUSTIFY}
        for c in json.loads(content):
            t=c["type"]
            if t=="heading":
                doc.add_heading(c["text"],level=c.get("level",1))
            elif t=="para":
                p=doc.add_paragraph()
                r=p.add_run(c["text"])
                if c.get("bold"):
                    r.bold=True
                if c.get("italic"):
                    r.italic=True
                if c.get("underline"):
                    r.underline=True
                if c.get("size"):
                    r.font.size=Pt(c["size"])
                if c.get("color"):
                    r.font.color.rgb=RGBColor.from_string(c["color"])
                if c.get("align"):
                    p.alignment=aligns.get(c["align"],WD_ALIGN_PARAGRAPH.LEFT)
            elif t=="table":
                rows=c["rows"]
                if not rows:
                    continue
                tb=doc.add_table(len(rows),len(rows[0]),style='Table Grid')
                for i,row in enumerate(rows):
                    for j,v in enumerate(row):
                        tb.cell(i,j).text=str(v)
                if c.get("header") and rows:
                    for cell in tb.rows[0].cells:
                        for p in cell.paragraphs:
                            for r in p.runs:
                                r.bold=True
                if c.get("col_widths"):
                    for i,w in enumerate(c["col_widths"]):
                        for row in tb.rows:
                            row.cells[i].width=Cm(w)
            elif t=="image":
                doc.add_picture(str(R(c["path"])),width=Inches(c.get("width",5)))
            elif t=="list":
                sty='List Bullet' if c.get("style","bullet")=="bullet" else 'List Number'
                for item in c.get("items",[]):
                    doc.add_paragraph(item,style=sty)
            elif t=="hr":
                p=doc.add_paragraph()
                p.add_run().add_break()
                from docx.oxml.ns import qn
                pPr=p._p.get_or_add_pPr()
                pBdr=pPr.makeelement(qn('w:pBdr'),{})
                bottom=pBdr.makeelement(qn('w:bottom'),{qn('w:val'):'single',qn('w:sz'):'6',qn('w:space'):'1',qn('w:color'):'auto'})
                pBdr.append(bottom)
                pPr.append(pBdr)
            elif t=="toc":
                doc.add_paragraph("[TOC - 在Word中按Ctrl+A后F9更新]")
            elif t=="break":
                doc.add_page_break()
        p=str(R(path))
        doc.save(p)
        return J(path=p)
    except Exception as e:
        return J(False,err=str(e))
@mcp.tool(name="docx_read")
async def docx_read(path:str)->str:
    """读取Word文档全部内容(段落+表格+图片数)"""
    from docx import Document
    try:
        doc=Document(str(R(path)))
        paras=[{"text":p.text,"style":p.style.name,"bold":any(r.bold for r in p.runs if r.bold)} for p in doc.paragraphs if p.text.strip()]
        tables=[]
        for tb in doc.tables:
            tables.append([[c.text for c in row.cells] for row in tb.rows])
        imgs=sum(1 for rel in doc.part.rels.values() if "image" in rel.reltype)
        return J(paragraphs=paras,tables=tables,images=imgs,
                 sections=len(doc.sections),total_paras=len(doc.paragraphs))
    except Exception as e:
        return J(False,err=str(e))
@mcp.tool(name="docx_replace")
async def docx_replace(path:str,replacements:str="{}",output:str="")->str:
    """Word文档查找替换
    replacements: JSON对象 {"旧文本":"新文本","{{name}}":"张三"}
    output: 输出路径,空则覆盖原文件"""
    from docx import Document
    try:
        p=str(R(path))
        doc=Document(p)
        reps=json.loads(replacements) if isinstance(replacements,str) else replacements
        for para in doc.paragraphs:
            for old,new in reps.items():
                if old in para.text:
                    for run in para.runs:
                        if old in run.text:
                            run.text=run.text.replace(old,new)
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    for para in cell.paragraphs:
                        for old,new in reps.items():
                            if old in para.text:
                                for run in para.runs:
                                    if old in run.text:
                                        run.text=run.text.replace(old,new)
        op=str(R(output)) if output else p
        doc.save(op)
        return J(path=op,replaced=len(reps))
    except Exception as e:
        return J(False,err=str(e))
# --- XLSX SUBSYSTEM: Workbook operations and charting ---
@mcp.tool(name="xlsx_create")
async def xlsx_create(path:str,sheets:str='[{"name":"Sheet1","data":[]}]')->str:
    """创建Excel
    sheets: JSON数组:
      {"name":"Sheet1","data":[[...],[...]],"widths":[15,20],
       "freeze":"A2","auto_filter":true,
       "merge":[{"range":"A1:C1","value":"标题"}],
       "styles":[{"range":"A1:Z1","bold":true,"bg":"FFFF00","font_color":"000000"}]}"""
    from openpyxl import Workbook
    from openpyxl.utils import get_column_letter
    from openpyxl.styles import Font,PatternFill,Alignment
    try:
        wb=Workbook()
        wb.remove(wb.active)
        for s in json.loads(sheets):
            ws=wb.create_sheet(s.get("name","Sheet"))
            for row in s.get("data",[]):
                ws.append(row)
            for i,w in enumerate(s.get("widths",[]),1):
                ws.column_dimensions[get_column_letter(i)].width=w
            if s.get("freeze"):
                ws.freeze_panes=s["freeze"]
            if s.get("auto_filter") and ws.max_row:
                ws.auto_filter.ref=ws.dimensions
            for m in s.get("merge",[]):
                ws.merge_cells(m["range"])
                if m.get("value"):
                    ws[m["range"].split(":")[0]]=m["value"]
            for st in s.get("styles",[]):
                font=Font(bold=st.get("bold",False),
                         color=st.get("font_color","000000"),
                         size=st.get("font_size",11))
                fill=PatternFill(start_color=st.get("bg","FFFFFF"),
                                end_color=st.get("bg","FFFFFF"),
                                fill_type="solid") if st.get("bg") else None
                for row in ws[st["range"]]:
                    for cell in (row if isinstance(row,tuple) else [row]):
                        cell.font=font
                        if fill:
                            cell.fill=fill
        p=str(R(path))
        wb.save(p)
        return J(path=p)
    except Exception as e:
        return J(False,err=str(e))
@mcp.tool(name="xlsx_read")
async def xlsx_read(path:str,sheet:str="",cell_range:str="")->str:
    """读取Excel。sheet空=第一个。cell_range如'A1:C10'"""
    from openpyxl import load_workbook
    try:
        wb=load_workbook(str(R(path)),data_only=True)
        ws=wb[sheet] if sheet else wb.active
        if cell_range:
            data=[[c.value for c in row] for row in ws[cell_range]]
        else:
            data=[[c.value for c in row] for row in ws.iter_rows()]
        return J(sheet=ws.title,data=data,rows=len(data),
                 cols=ws.max_column,sheets=wb.sheetnames)
    except Exception as e:
        return J(False,err=str(e))
@mcp.tool(name="xlsx_write")
async def xlsx_write(path:str,sheet:str="",writes:str="[]")->str:
    """写入Excel单元格
    writes: JSON数组:
      {"cell":"A1","value":"hello"}
      {"cell":"A1","formula":"=SUM(B1:B10)"}
      {"row":1,"col":1,"value":42}
      {"range":"A1:A10","values":[1,2,3,4,5,6,7,8,9,10]}"""
    from openpyxl import load_workbook
    try:
        p=str(R(path))
        wb=load_workbook(p)
        ws=wb[sheet] if sheet else wb.active
        for w in json.loads(writes):
            if "formula" in w:
                ws[w["cell"]]=w["formula"]
            elif "range" in w and "values" in w:
                from openpyxl.utils import range_boundaries
                mc,mr,xc,xr=range_boundaries(w["range"])
                vals=w["values"]
                idx=0
                for r in range(mr,xr+1):
                    for c in range(mc,xc+1):
                        if idx<len(vals):
                            ws.cell(row=r,column=c,value=vals[idx])
                            idx+=1
            elif "cell" in w:
                ws[w["cell"]]=w["value"]
            else:
                ws.cell(row=w["row"],column=w["col"],value=w["value"])
        wb.save(p)
        return J(path=p)
    except Exception as e:
        return J(False,err=str(e))
@mcp.tool(name="xlsx_chart")
async def xlsx_chart(path:str,sheet:str="",chart_config:str="{}")->str:
    """在Excel中插入图表
    chart_config: {"type":"bar"|"line"|"pie"|"scatter"|"area",
      "title":"图表标题","data_range":"A1:B10","categories_range":"A1:A10",
      "position":"D1","width":15,"height":10}"""
    from openpyxl import load_workbook
    from openpyxl.chart import BarChart,LineChart,PieChart,ScatterChart,AreaChart,Reference
    try:
        p=str(R(path))
        wb=load_workbook(p)
        ws=wb[sheet] if sheet else wb.active
        cfg=json.loads(chart_config) if isinstance(chart_config,str) else chart_config
        ct=cfg.get("type","bar")
        charts={"bar":BarChart,"line":LineChart,"pie":PieChart,"scatter":ScatterChart,"area":AreaChart}
        chart=charts.get(ct,BarChart)()
        chart.title=cfg.get("title","")
        chart.width=cfg.get("width",15)
        chart.height=cfg.get("height",10)
        dr=cfg.get("data_range","B1:B10")
        from openpyxl.utils import range_boundaries
        mc,mr,xc,xr=range_boundaries(dr)
        data=Reference(ws,min_col=mc,min_row=mr,max_col=xc,max_row=xr)
        chart.add_data(data,titles_from_data=True)
        if cfg.get("categories_range"):
            cc,cr,xcc,xcr=range_boundaries(cfg["categories_range"])
            cats=Reference(ws,min_col=cc,min_row=cr,max_col=xcc,max_row=xcr)
            chart.set_categories(cats)
        ws.add_chart(chart,cfg.get("position","D1"))
        wb.save(p)
        return J(path=p)
    except Exception as e:
        return J(False,err=str(e))
# --- PDF SUBSYSTEM: Generation, parsing, and post-processing ---
@mcp.tool(name="pdf_create")
async def pdf_create(path:str,content:str="[]",page_size:str="A4")->str:
    """创建PDF
    content: JSON数组:
      {"type":"title","text":"","size":24}
      {"type":"text","text":"","size":12,"bold":false,"color":"#000000"}
      {"type":"image","path":"","w":400,"h":300}
      {"type":"table","data":[["a","b"],["c","d"]],"col_widths":[200,200]}
      {"type":"spacer","h":20}
      {"type":"line","x1":50,"y1":0,"x2":550,"y2":0,"color":"#000000","width":1}
      {"type":"break"}"""
    from reportlab.lib.pagesizes import A4,letter
    from reportlab.platypus import SimpleDocTemplate,Paragraph,Spacer,Image,Table,TableStyle,PageBreak
    from reportlab.lib.styles import getSampleStyleSheet,ParagraphStyle
    from reportlab.lib import colors
    from reportlab.lib.units import mm
    try:
        # 注册中文字体(优先微软雅黑，Unicode覆盖最全)
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.ttfonts import TTFont
        _cn_font="Helvetica"
        _cn_font_bd="Helvetica-Bold"
        _fonts_dir=os.path.join(os.environ.get("WINDIR","C:\\Windows"),"Fonts")
        _font_candidates=[
            ("msyh.ttc","MSYH","msyhbd.ttc","MSYHBD"),
            ("simhei.ttf","SimHei","simhei.ttf","SimHei"),
            ("simsun.ttc","SimSun","simsun.ttc","SimSun"),
        ]
        for _reg,_alias,_reg_bd,_alias_bd in _font_candidates:
            _fp=os.path.join(_fonts_dir,_reg)
            _fp_bd=os.path.join(_fonts_dir,_reg_bd)
            if os.path.isfile(_fp):
                try:
                    pdfmetrics.registerFont(TTFont(_alias,_fp))
                    _cn_font=_alias
                    if os.path.isfile(_fp_bd) and _reg_bd!=_reg:
                        pdfmetrics.registerFont(TTFont(_alias_bd,_fp_bd))
                        _cn_font_bd=_alias_bd
                    else:
                        _cn_font_bd=_alias
                    break
                except:
                    pass
        ps={"A4":A4,"letter":letter}.get(page_size,A4)
        p=str(R(path))
        doc=SimpleDocTemplate(p,pagesize=ps,
            topMargin=50,bottomMargin=50,leftMargin=55,rightMargin=55)
        sty=getSampleStyleSheet()
        story=[]
        for c in json.loads(content):
            t=c["type"]
            if t in("text","title"):
                if t=="title":
                    s=ParagraphStyle(f"_t{id(c)}",fontName=_cn_font_bd,
                        fontSize=c.get("size",22),leading=c.get("size",22)*1.4,
                        spaceBefore=16,spaceAfter=12,alignment=1)
                else:
                    sz=c.get("size",11)
                    s=ParagraphStyle(f"_n{id(c)}",fontName=_cn_font,
                        fontSize=sz,leading=sz*1.6,
                        spaceBefore=2,spaceAfter=4)
                if c.get("bold"):
                    s.fontName=_cn_font_bd
                if c.get("color"):
                    col=c["color"].lstrip("#")
                    s.textColor=colors.HexColor(f"#{col}")
                if c.get("align"):
                    from reportlab.lib.enums import TA_CENTER,TA_RIGHT,TA_JUSTIFY
                    aligns={"center":TA_CENTER,"right":TA_RIGHT,"justify":TA_JUSTIFY}
                    s.alignment=aligns.get(c["align"],0)
                # heading类型:标题样式带下划线装饰
                if c.get("heading"):
                    lvl=c.get("heading",1)
                    sz=[18,15,13][min(lvl,3)-1]
                    s.fontSize=sz
                    s.leading=sz*1.5
                    s.spaceBefore=14
                    s.spaceAfter=6
                    s.fontName=_cn_font_bd
                    s.textColor=colors.HexColor("#1a1a2e")
                story.append(Paragraph(c["text"],s))
            elif t=="heading":
                lvl=c.get("level",1)
                sz=[20,16,13][min(lvl,3)-1]
                s=ParagraphStyle(f"_h{id(c)}",fontName=_cn_font_bd,
                    fontSize=sz,leading=sz*1.5,
                    spaceBefore=16,spaceAfter=8,
                    textColor=colors.HexColor("#1a1a2e"))
                story.append(Paragraph(c["text"],s))
            elif t=="image":
                story.append(Image(str(R(c["path"])),c.get("w",400),c.get("h",300)))
            elif t=="table":
                cw=c.get("col_widths")
                tb=Table(c["data"],colWidths=cw)
                ts=[
                    ('GRID',(0,0),(-1,-1),0.5,colors.Color(0.4,0.4,0.4)),
                    ('BACKGROUND',(0,0),(-1,0),colors.HexColor("#2c3e50")),
                    ('TEXTCOLOR',(0,0),(-1,0),colors.white),
                    ('FONTNAME',(0,0),(-1,-1),_cn_font),
                    ('FONTNAME',(0,0),(-1,0),_cn_font_bd),
                    ('FONTSIZE',(0,0),(-1,-1),9),
                    ('FONTSIZE',(0,0),(-1,0),10),
                    ('ALIGN',(0,0),(-1,-1),'CENTER'),
                    ('VALIGN',(0,0),(-1,-1),'MIDDLE'),
                    ('TOPPADDING',(0,0),(-1,-1),5),
                    ('BOTTOMPADDING',(0,0),(-1,-1),5),
                    ('ROWBACKGROUNDS',(0,1),(-1,-1),[colors.white,colors.HexColor("#f8f9fa")]),
                ]
                tb.setStyle(TableStyle(ts))
                story.append(Spacer(1,6))
                story.append(tb)
                story.append(Spacer(1,8))
            elif t=="spacer":
                story.append(Spacer(1,c.get("h",20)))
            elif t=="break":
                story.append(PageBreak())
        doc.build(story)
        return J(path=p)
    except Exception as e:
        return J(False,err=str(e))
@mcp.tool(name="pdf_read")
async def pdf_read(path:str,pages:str="")->str:
    """读取PDF文本。pages如'0,1,2'(从0开始),空=全部"""
    import fitz
    try:
        doc=fitz.open(str(R(path)))
        total=len(doc)
        ps=[int(x) for x in pages.split(",") if x.strip()] if pages else list(range(total))
        out=[{"page":i,"text":doc[i].get_text()} for i in ps if i<total]
        doc.close()
        return J(pages=out,total=total)
    except Exception as e:
        return J(False,err=str(e))
@mcp.tool(name="pdf_merge")
async def pdf_merge(files:str="[]",output:str="merged.pdf")->str:
    """合并多个PDF"""
    import fitz
    try:
        flist=json.loads(files)
        if not flist:
            return J(False,err="files数组为空")
        out=fitz.open()
        for f in flist:
            fp=str(R(f))
            src=fitz.open(fp)
            out.insert_pdf(src)
            src.close()
        if len(out)==0:
            return J(False,err="合并后文档无页面")
        p=str(R(output))
        out.save(p)
        out.close()
        return J(path=p)
    except Exception as e:
        return J(False,err=str(e))
@mcp.tool(name="pdf_split")
async def pdf_split(path:str,page_ranges:str="",output_dir:str="")->str:
    """拆分PDF
    page_ranges: "0-2,3-5,6-10" 按范围拆分, 空=每页一个文件
    output_dir: 输出目录"""
    import fitz
    try:
        doc=fitz.open(str(R(path)))
        od=Path(output_dir) if output_dir else WD
        od.mkdir(parents=True,exist_ok=True)
        outs=[]
        if page_ranges:
            for rng in page_ranges.split(","):
                parts=rng.strip().split("-")
                s,e=int(parts[0]),int(parts[-1])
                nd=fitz.open()
                for i in range(s,min(e+1,len(doc))):
                    nd.insert_pdf(doc,from_page=i,to_page=i)
                op=str(od/f"split_{s}-{e}.pdf")
                nd.save(op)
                nd.close()
                outs.append(op)
        else:
            for i in range(len(doc)):
                nd=fitz.open()
                nd.insert_pdf(doc,from_page=i,to_page=i)
                op=str(od/f"page_{i+1}.pdf")
                nd.save(op)
                nd.close()
                outs.append(op)
        doc.close()
        return J(files=outs,count=len(outs))
    except Exception as e:
        return J(False,err=str(e))
@mcp.tool(name="pdf_watermark")
async def pdf_watermark(path:str,text:str="WATERMARK",output:str="",
                        font_size:int=50,color:str="0.8 0.8 0.8",rotation:int=45)->str:
    """给PDF添加文字水印"""
    import fitz
    try:
        doc=fitz.open(str(R(path)))
        for page in doc:
            r=page.rect
            cx,cy=r.width/2,r.height/2
            page.insert_text((cx-100,cy),text,fontsize=font_size,
                           color=[float(x) for x in color.split()],rotate=rotation)
        op=str(R(output)) if output else str(R(path))
        doc.save(op)
        doc.close()
        return J(path=op)
    except Exception as e:
        return J(False,err=str(e))
@mcp.tool(name="pdf_to_images")
async def pdf_to_images(path:str,output_dir:str="",dpi:int=200,fmt:str="png")->str:
    """PDF转图片"""
    import fitz
    try:
        doc=fitz.open(str(R(path)))
        od=Path(output_dir) if output_dir else WD
        od.mkdir(parents=True,exist_ok=True)
        outs=[]
        mat=fitz.Matrix(dpi/72,dpi/72)
        for i,page in enumerate(doc):
            pix=page.get_pixmap(matrix=mat)
            op=str(od/f"page_{i+1}.{fmt}")
            pix.save(op)
            outs.append(op)
        doc.close()
        return J(files=outs,count=len(outs))
    except Exception as e:
        return J(False,err=str(e))
# --- IMAGE SUBSYSTEM: Raster processing pipeline ---
@mcp.tool(name="img_process")
async def img_process(src:str,dst:str="",ops:str="[]")->str:
    """图像处理(增强版)
    ops: JSON数组,按序执行:
      {"op":"resize","w":800,"h":600}  {"op":"resize_ratio","ratio":0.5}
      {"op":"crop","l":0,"t":0,"r":100,"b":100}
      {"op":"rotate","angle":90}  {"op":"flip","dir":"h"|"v"}
      {"op":"blur","r":5}  {"op":"sharpen"}  {"op":"gray"}
      {"op":"brightness","f":1.5}  {"op":"contrast","f":1.5}  {"op":"saturation","f":1.5}
      {"op":"text","text":"水印","x":10,"y":10,"size":36,"color":"white","opacity":128}
      {"op":"border","size":10,"color":"black"}
      {"op":"round_corners","radius":20}
      {"op":"convert","mode":"RGBA"|"RGB"|"L"|"CMYK"}
      {"op":"thumbnail","size":256}
      {"op":"overlay","path":"logo.png","x":10,"y":10,"opacity":0.5}
    dst空则覆盖src"""
    from PIL import Image,ImageFilter,ImageEnhance,ImageDraw,ImageFont,ImageOps
    try:
        img=Image.open(str(R(src)))
        if img.mode=="P":
            img=img.convert("RGBA")
        for o in json.loads(ops):
            op=o["op"]
            if op=="resize":
                img=img.resize((o["w"],o["h"]),Image.LANCZOS)
            elif op=="resize_ratio":
                img=img.resize((int(img.width*o["ratio"]),int(img.height*o["ratio"])),Image.LANCZOS)
            elif op=="crop":
                img=img.crop((o["l"],o["t"],o["r"],o["b"]))
            elif op=="rotate":
                img=img.rotate(o["angle"],expand=True,resample=Image.BICUBIC)
            elif op=="flip":
                img=img.transpose(Image.FLIP_LEFT_RIGHT if o["dir"]=="h" else Image.FLIP_TOP_BOTTOM)
            elif op=="blur":
                img=img.filter(ImageFilter.GaussianBlur(o.get("r",5)))
            elif op=="sharpen":
                img=img.filter(ImageFilter.SHARPEN)
            elif op=="gray":
                img=img.convert("L")
            elif op=="brightness":
                img=ImageEnhance.Brightness(img).enhance(o["f"])
            elif op=="contrast":
                img=ImageEnhance.Contrast(img).enhance(o["f"])
            elif op=="saturation":
                img=ImageEnhance.Color(img).enhance(o["f"])
            elif op=="convert":
                img=img.convert(o["mode"])
            elif op=="thumbnail":
                s=o["size"]
                img.thumbnail((s,s),Image.LANCZOS)
            elif op=="border":
                img=ImageOps.expand(img,border=o.get("size",10),fill=o.get("color","black"))
            elif op=="overlay":
                overlay=Image.open(str(R(o["path"])))
                if o.get("opacity",1)<1:
                    overlay=overlay.convert("RGBA")
                    alpha=overlay.split()[3].point(lambda p:int(p*o["opacity"]))
                    overlay.putalpha(alpha)
                img.paste(overlay,(o.get("x",0),o.get("y",0)),overlay if overlay.mode=="RGBA" else None)
            elif op=="text":
                if img.mode!="RGBA":
                    img=img.convert("RGBA")
                txt_layer=Image.new("RGBA",img.size,(0,0,0,0))
                d=ImageDraw.Draw(txt_layer)
                try:
                    ft=ImageFont.truetype("arial.ttf",o.get("size",36))
                except:
                    ft=ImageFont.load_default()
                fill=o.get("color","white")
                d.text((o.get("x",10),o.get("y",10)),o["text"],fill=fill,font=ft)
                img=Image.alpha_composite(img,txt_layer)
        p=str(R(dst or src))
        img.save(p)
        return J(path=p,size=list(img.size),mode=img.mode)
    except Exception as e:
        return J(False,err=str(e))
@mcp.tool(name="img_create")
async def img_create(path:str,w:int=800,h:int=600,color:str="white")->str:
    """创建纯色图像"""
    from PIL import Image
    try:
        img=Image.new("RGB",(w,h),color)
        p=str(R(path))
        img.save(p)
        return J(path=p)
    except Exception as e:
        return J(False,err=str(e))
@mcp.tool(name="img_info")
async def img_info(path:str)->str:
    """获取图像详细信息"""
    from PIL import Image
    try:
        img=Image.open(str(R(path)))
        info={"size":list(img.size),"mode":img.mode,"format":img.format,
              "megapixels":round(img.width*img.height/1e6,2)}
        if img.info.get("dpi"):
            info["dpi"]=img.info["dpi"]
        try:
            exif=img._getexif()
            if exif:
                info["has_exif"]=True
        except:
            pass
        return J(**info)
    except Exception as e:
        return J(False,err=str(e))
@mcp.tool(name="img_convert")
async def img_convert(src:str,dst:str,quality:int=95)->str:
    """图像格式转换(png/jpg/bmp/webp/tiff/gif)"""
    from PIL import Image
    try:
        img=Image.open(str(R(src)))
        if dst.lower().endswith(('.jpg','.jpeg')) and img.mode in ('RGBA','P'):
            img=img.convert('RGB')
        p=str(R(dst))
        img.save(p,quality=quality)
        return J(path=p,size=list(img.size))
    except Exception as e:
        return J(False,err=str(e))
@mcp.tool(name="img_composite")
async def img_composite(images:str="[]",output:str="composite.png",
                        direction:str="horizontal",gap:int=0,bg:str="white")->str:
    """拼接多张图片
    images: JSON数组 ["img1.png","img2.png",...]
    direction: horizontal|vertical|grid
    gap: 图片间距(px)"""
    from PIL import Image
    try:
        imgs=[Image.open(str(R(p))) for p in json.loads(images)]
        if not imgs:
            return J(False,err="no images")
        if direction=="horizontal":
            h=max(i.height for i in imgs)
            w=sum(i.width for i in imgs)+gap*(len(imgs)-1)
            out=Image.new("RGB",(w,h),bg)
            x=0
            for i in imgs:
                out.paste(i,(x,0))
                x+=i.width+gap
        elif direction=="vertical":
            w=max(i.width for i in imgs)
            h=sum(i.height for i in imgs)+gap*(len(imgs)-1)
            out=Image.new("RGB",(w,h),bg)
            y=0
            for i in imgs:
                out.paste(i,(0,y))
                y+=i.height+gap
        else: # grid
            n=len(imgs)
            cols=int(n**0.5)+1
            rows=(n+cols-1)//cols
            mw=max(i.width for i in imgs)
            mh=max(i.height for i in imgs)
            out=Image.new("RGB",(cols*(mw+gap)-gap,rows*(mh+gap)-gap),bg)
            for idx,i in enumerate(imgs):
                r,c=divmod(idx,cols)
                out.paste(i,(c*(mw+gap),r*(mh+gap)))
        p=str(R(output))
        out.save(p)
        return J(path=p,size=list(out.size))
    except Exception as e:
        return J(False,err=str(e))
# --- BLENDER SUBSYSTEM: Headless modeling and rendering ---
@mcp.tool(name="blender_exec")
async def blender_exec(script:str,blend_file:str="",timeout:int=300)->str:
    """在Blender后台执行Python脚本
    script: bpy Python代码
    blend_file: 可选,.blend文件路径
    常用: bpy.ops.mesh.primitive_xxx_add / bpy.ops.render.render(write_still=True)"""
    try:
        sf=WD/"_bpy.py"
        sf.write_text(script,encoding="utf-8")
        cmd=[BLENDER,"--background"]
        if blend_file:
            cmd.append(str(R(blend_file)))
        cmd.extend(["--python",str(sf)])
        o,e,c=_run(cmd,timeout=timeout)
        return J(stdout=o,stderr=e,code=c)
    except subprocess.TimeoutExpired:
        return J(False,err=f"超时({timeout}s)")
    except Exception as e:
        return J(False,err=str(e))
@mcp.tool(name="blender_render")
async def blender_render(blend_file:str,output:str,engine:str="CYCLES",
                         rx:int=1920,ry:int=1080,samples:int=128,frame:int=1)->str:
    """渲染Blender场景
    engine: CYCLES | BLENDER_EEVEE_NEXT
    output: 输出图片路径(.png/.jpg/.exr)"""
    op=str(R(output)).replace("\\","\\\\")
    script=f"""import bpy
s=bpy.context.scene
s.render.engine='{engine}'
s.render.resolution_x={rx}
s.render.resolution_y={ry}
{'s.cycles.samples='+str(samples) if engine=='CYCLES' else ''}
s.render.filepath=r'{op}'
s.frame_set({frame})
bpy.ops.render.render(write_still=True)
print('RENDER_DONE:',s.render.filepath)
"""
    try:
        cmd=[BLENDER,"--background",str(R(blend_file)),"--python-expr",script]
        o,e,c=_run(cmd,timeout=600)
        return J(path=str(R(output)),code=c)
    except Exception as e:
        return J(False,err=str(e))
@mcp.tool(name="blender_scene")
async def blender_scene(action:str,kw:str="{}",blend_file:str="")->str:
    """快捷Blender场景操作
    action: 
      "add_mesh" - 添加网格 (mesh="cube"|"sphere"|"cylinder"|"cone"|"torus"|"monkey", loc=[0,0,0], scale=[1,1,1])
      "add_light" - 添加灯光 (light="POINT"|"SUN"|"SPOT"|"AREA", loc=[0,0,5], energy=1000)
      "add_camera" - 添加相机 (loc=[7,−6,5], rot=[1.1,0,0.8])
      "set_material" - 设置材质 (obj="Cube", color=[0.8,0.1,0.1,1])
      "delete" - 删除对象 (obj="Cube")
      "export" - 导出 (fmt="fbx"|"obj"|"stl"|"gltf", path="out.fbx")
      "save" - 保存blend文件 (path="scene.blend")
    kw: JSON对象,传递action所需参数
    blend_file: 可选,在此文件基础上操作"""
    kw=json.loads(kw) if isinstance(kw,str) else kw
    meshes={"cube":"primitive_cube_add","sphere":"primitive_uv_sphere_add",
            "cylinder":"primitive_cylinder_add","cone":"primitive_cone_add",
            "torus":"primitive_torus_add","monkey":"primitive_monkey_add",
            "plane":"primitive_plane_add","ico":"primitive_ico_sphere_add"}
    exports={"fbx":"export_scene.fbx","obj":"export_scene.obj",
             "stl":"export_mesh.stl","gltf":"export_scene.gltf"}
    lines=["import bpy","import math"]
    a=action
    if a=="add_mesh":
        m=kw.get("mesh","cube")
        fn=meshes.get(m,f"primitive_{m}_add")
        loc=kw.get("loc",[0,0,0])
        scl=kw.get("scale",[1,1,1])
        lines.append(f"bpy.ops.mesh.{fn}(location={loc},scale={scl})")
    elif a=="add_light":
        lt=kw.get("light","POINT")
        loc=kw.get("loc",[0,0,5])
        e=kw.get("energy",1000)
        lines.append(f"bpy.ops.object.light_add(type='{lt}',location={loc})")
        lines.append(f"bpy.context.object.data.energy={e}")
    elif a=="add_camera":
        loc=kw.get("loc",[7,-6,5])
        rot=kw.get("rot",[1.1,0,0.8])
        lines.append(f"bpy.ops.object.camera_add(location={loc},rotation={rot})")
        lines.append("bpy.context.scene.camera=bpy.context.object")
    elif a=="set_material":
        obj=kw.get("obj","Cube")
        color=kw.get("color",[0.8,0.1,0.1,1])
        lines.append(f"o=bpy.data.objects['{obj}']")
        lines.append(f"mat=bpy.data.materials.new('Mat'); mat.use_nodes=True")
        lines.append(f"mat.node_tree.nodes['Principled BSDF'].inputs['Base Color'].default_value={color}")
        lines.append(f"o.data.materials.append(mat)")
    elif a=="delete":
        obj=kw.get("obj","Cube")
        lines.append(f"bpy.data.objects.remove(bpy.data.objects['{obj}'],do_unlink=True)")
    elif a=="export":
        fmt=kw.get("fmt","fbx")
        ep=str(R(kw.get("path",f"out.{fmt}"))).replace("\\","\\\\")
        fn=exports.get(fmt,f"export_scene.{fmt}")
        lines.append(f"bpy.ops.{fn}(filepath=r'{ep}')")
    elif a=="save":
        sp=str(R(kw.get("path","scene.blend"))).replace("\\","\\\\")
        lines.append(f"bpy.ops.wm.save_as_mainfile(filepath=r'{sp}')")
    lines.append("print('DONE')")
    return await blender_exec("\n".join(lines),blend_file)
# --- SVG SUBSYSTEM: Vector description synthesis ---
@mcp.tool(name="svg_create")
async def svg_create(path:str,w:int=800,h:int=600,elements:str="[]",bg:str="")->str:
    """创建SVG
    elements: JSON数组:
      {"tag":"rect","x":0,"y":0,"w":100,"h":50,"fill":"blue","stroke":"black","rx":5}
      {"tag":"circle","cx":50,"cy":50,"r":30,"fill":"red","opacity":0.8}
      {"tag":"line","x1":0,"y1":0,"x2":100,"y2":100,"stroke":"black","stroke-width":2}
      {"tag":"text","x":10,"y":30,"text":"Hello","font-size":16,"fill":"black","font-family":"Arial"}
      {"tag":"ellipse","cx":50,"cy":50,"rx":40,"ry":20,"fill":"green"}
      {"tag":"path","d":"M10 10 L90 90","stroke":"black","fill":"none"}
      {"tag":"polygon","points":"0,0 50,50 100,0","fill":"yellow"}
      {"tag":"polyline","points":"0,0 50,50 100,0","stroke":"blue","fill":"none"}
      {"tag":"group","transform":"translate(100,100) rotate(45)","children":[...]}
      {"tag":"defs","gradient":{"id":"g1","type":"linear","stops":[{"offset":"0%","color":"red"},{"offset":"100%","color":"blue"}]}}
    bg: 背景色,空=透明"""
    try:
        parts=[f'<svg xmlns="http://www.w3.org/2000/svg" width="{w}" height="{h}" viewBox="0 0 {w} {h}">']
        if bg:
            parts.append(f'  <rect width="{w}" height="{h}" fill="{bg}"/>')
        def render_el(e,indent=2):
            tag=e.get("tag","rect")
            sp=" "*indent
            if tag=="group":
                tr=e.get("transform","")
                parts.append(f'{sp}<g transform="{tr}">')
                for ch in e.get("children",[]):
                    render_el(ch,indent+2)
                parts.append(f'{sp}</g>')
            elif tag=="defs" and e.get("gradient"):
                g=e["gradient"]
                gt=g.get("type","linear")
                parts.append(f'{sp}<defs>')
                parts.append(f'{sp}  <{gt}Gradient id="{g["id"]}">')
                for s in g.get("stops",[]):
                    parts.append(f'{sp}    <stop offset="{s["offset"]}" stop-color="{s["color"]}"/>')
                parts.append(f'{sp}  </{gt}Gradient>')
                parts.append(f'{sp}</defs>')
            elif tag=="text":
                txt=e.pop("text","")
                e.pop("tag",None)
                attrs=" ".join(f'{k}="{v}"' for k,v in e.items())
                parts.append(f"{sp}<text {attrs}>{txt}</text>")
            else:
                e2=dict(e)
                e2.pop("tag",None)
                if tag=="rect":
                    e2.setdefault("width",e2.pop("w",100))
                    e2.setdefault("height",e2.pop("h",50))
                attrs=" ".join(f'{k}="{v}"' for k,v in e2.items())
                parts.append(f"{sp}<{tag} {attrs}/>")
        for e in json.loads(elements):
            render_el(e)
        parts.append("</svg>")
        p=str(R(path))
        Path(p).write_text("\n".join(parts),encoding="utf-8")
        return J(path=p)
    except Exception as e:
        return J(False,err=str(e))
# --- CHART SUBSYSTEM: Statistical plotting utilities ---
@mcp.tool(name="chart_create")
async def chart_create(path:str,chart_type:str="line",title:str="",
                       x_label:str="",y_label:str="",
                       datasets:str='[{"label":"data","x":[1,2,3],"y":[4,5,6]}]',
                       w:float=10,h:float=6,style:str="default",dpi:int=150,
                       grid:bool=True,legend_pos:str="best",
                       annotations:str="[]",x_rotation:int=0)->str:
    """创建图表并保存为图片(增强版)
    chart_type: line|bar|scatter|pie|hist|area|box|hbar|stackbar|heatmap|radar|stem|step
    datasets: JSON数组: {"label":"","x":[],"y":[],"color":"","marker":"o","linestyle":"-","linewidth":2}
    pie: {"labels":[],"values":[],"colors":[],"explode":[]}
    hist: {"y":[],"bins":20}
    heatmap: {"data":[[...]],"xlabels":[],"ylabels":[],"cmap":"viridis"}
    annotations: [{"text":"标注","x":2,"y":5,"arrow":true}]
    legend_pos: best|upper right|upper left|lower right|lower left|center"""
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    import numpy as np
    try:
        if style in plt.style.available:
            plt.style.use(style)
        fig,ax=plt.subplots(figsize=(w,h))
        ds=json.loads(datasets)
        ct=chart_type
        if ct=="pie":
            d=ds[0]
            ax.pie(d["values"],labels=d["labels"],colors=d.get("colors"),
                            explode=d.get("explode"),autopct='%1.1f%%',startangle=90)
        elif ct=="heatmap":
            d=ds[0]
            data=np.array(d["data"])
            im=ax.imshow(data,cmap=d.get("cmap","viridis"),aspect="auto")
            fig.colorbar(im,ax=ax)
            if d.get("xlabels"):
                ax.set_xticks(range(len(d["xlabels"])))
                ax.set_xticklabels(d["xlabels"])
            if d.get("ylabels"):
                ax.set_yticks(range(len(d["ylabels"])))
                ax.set_yticklabels(d["ylabels"])
        elif ct=="radar":
            d=ds[0]
            labels=d.get("labels",[])
            vals=d["y"]
            N=len(labels)
            angles=np.linspace(0,2*np.pi,N,endpoint=False).tolist()
            vals=vals+vals[:1]
            angles+=angles[:1]
            ax=fig.add_subplot(111,polar=True)
            ax.plot(angles,vals)
            ax.fill(angles,vals,alpha=0.25)
            ax.set_xticks(angles[:-1])
            ax.set_xticklabels(labels)
        else:
            bottom_stack=None
            for d in ds:
                lb=d.get("label","")
                cl=d.get("color")
                kw={"label":lb}
                if cl:
                    kw["color"]=cl
                if d.get("marker"):
                    kw["marker"]=d["marker"]
                if d.get("linestyle"):
                    kw["linestyle"]=d["linestyle"]
                if d.get("linewidth"):
                    kw["linewidth"]=d["linewidth"]
                x,y=d.get("x"),d.get("y",[])
                if ct=="line":
                    ax.plot(x,y,**kw)
                elif ct=="step":
                    ax.step(x,y,**kw)
                elif ct=="stem":
                    ax.stem(x,y,label=lb)
                elif ct=="bar":
                    ax.bar(x,y,**kw)
                elif ct=="hbar":
                    ax.barh(x,y,**kw)
                elif ct=="stackbar":
                    if bottom_stack is None:
                        ax.bar(x,y,**kw)
                        bottom_stack=np.array(y)
                    else:
                        ax.bar(x,y,bottom=bottom_stack,**kw)
                        bottom_stack+=np.array(y)
                elif ct=="scatter":
                    ax.scatter(x,y,s=d.get("size",20),**kw)
                elif ct=="area":
                    ax.fill_between(x,y,alpha=0.5,**kw)
                elif ct=="hist":
                    ax.hist(y,bins=d.get("bins",20),**kw)
                elif ct=="box":
                    ax.boxplot([dd["y"] for dd in ds],labels=[dd.get("label","") for dd in ds])
                    break
            if ct not in("box",):
                ax.legend(loc=legend_pos)
        # annotations
        for ann in json.loads(annotations):
            kw2={"fontsize":ann.get("fontsize",10)}
            if ann.get("arrow"):
                ax.annotate(ann["text"],xy=(ann["x"],ann["y"]),xytext=(ann.get("tx",ann["x"]+1),ann.get("ty",ann["y"]+1)),
                           arrowprops=dict(arrowstyle="->"),**kw2)
            else:
                ax.text(ann["x"],ann["y"],ann["text"],**kw2)
        if grid:
            ax.grid(True,alpha=0.3)
        if title:
            ax.set_title(title)
        if x_label:
            ax.set_xlabel(x_label)
        if y_label:
            ax.set_ylabel(y_label)
        if x_rotation:
            plt.xticks(rotation=x_rotation)
        plt.tight_layout()
        p=str(R(path))
        fig.savefig(p,dpi=dpi,bbox_inches='tight')
        plt.close()
        return J(path=p)
    except Exception as e:
        return J(False,err=str(e))
@mcp.tool(name="chart_subplot")
async def chart_subplot(path:str,rows:int=1,cols:int=2,
                        subplots:str="[]",w:float=14,h:float=6,dpi:int=150,
                        title:str="",style:str="")->str:
    """创建多子图
    subplots: JSON数组, 每项:
      {"title":"","chart_type":"line|bar|scatter|pie|area","datasets":[...],"x_label":"","y_label":"","grid":true}
    datasets格式同chart_create
    pie需要 x(标签) y(数值) colors(可选颜色列表)"""
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    try:
        if style:
            try:
                plt.style.use(style)
            except:
                pass
        fig,axes=plt.subplots(rows,cols,figsize=(w,h))
        if rows*cols==1:
            axes=[axes]
        elif rows==1 or cols==1:
            axes=axes.flatten()
        else:
            axes=axes.flatten()
        subs=json.loads(subplots)
        for i,sp in enumerate(subs):
            if i>=len(axes):
                break
            ax=axes[i]
            ct=sp.get("chart_type","line")
            if ct=="pie":
                d=sp.get("datasets",[{}])[0]
                labels=d.get("x",[])
                sizes=d.get("y",[])
                cs=d.get("colors",None)
                ax.pie(sizes,labels=labels,colors=cs,autopct='%1.1f%%',startangle=90,textprops={'fontsize':8})
                ax.set_aspect('equal')
            else:
                for d in sp.get("datasets",[]):
                    x,y=d.get("x"),d.get("y",[])
                    lb=d.get("label","")
                    kw={"label":lb}
                    if d.get("color"):
                        kw["color"]=d["color"]
                    if d.get("alpha"):
                        kw["alpha"]=d["alpha"]
                    if ct=="line":
                        ax.plot(x,y,**kw)
                    elif ct=="bar":
                        ax.bar(x,y,**kw)
                    elif ct=="scatter":
                        ax.scatter(x,y,**kw)
                    elif ct=="area":
                        ax.fill_between(x,y,alpha=d.get("alpha",0.4),**{k:v for k,v in kw.items() if k!="alpha"})
                if sp.get("x_label"):
                    ax.set_xlabel(sp["x_label"])
                if sp.get("y_label"):
                    ax.set_ylabel(sp["y_label"])
                if sp.get("grid",True):
                    ax.grid(True,alpha=0.3)
                ax.legend()
            if sp.get("title"):
                ax.set_title(sp["title"])
        if title:
            fig.suptitle(title,fontsize=14,fontweight='bold',y=1.02)
        plt.tight_layout()
        p=str(R(path))
        fig.savefig(p,dpi=dpi,bbox_inches='tight')
        plt.close()
        return J(path=p)
    except Exception as e:
        return J(False,err=str(e))
# --- MATLAB SUBSYSTEM: Batch script execution ---
@mcp.tool(name="matlab_exec")
async def matlab_exec(script:str,timeout:int=120)->str:
    """执行MATLAB脚本
    script: MATLAB代码(自动添加exit)
    常用: plot,surf,solve,eig,fft,ode45,simulink等
    输出图片用: saveas(gcf,'output.png')"""
    try:
        sf=WD/"omnirun.m"
        sf.write_text(script+"\nexit;\n",encoding="utf-8")
        mdir=str(WD).replace("\\","/")
        cmd=[MATLAB,"-batch",f"cd('{mdir}'); run('omnirun.m')"]
        o,e,c=_run(cmd,timeout=timeout)
        return J(stdout=o,stderr=e,code=c)
    except subprocess.TimeoutExpired:
        return J(False,err=f"MATLAB超时({timeout}s)")
    except Exception as e:
        return J(False,err=str(e))
@mcp.tool(name="matlab_eval")
async def matlab_eval(expr:str,timeout:int=60)->str:
    """快速执行MATLAB表达式并返回结果
    expr: 单行或多行MATLAB表达式
    例: "eig([1 2;3 4])" 或 "x=linspace(0,2*pi); y=sin(x); plot(x,y); saveas(gcf,'sin.png')" """
    try:
        cmd=[MATLAB,"-batch",expr]
        o,e,c=_run(cmd,timeout=timeout)
        return J(stdout=o,stderr=e,code=c)
    except Exception as e:
        return J(False,err=str(e))
# --- FFMPEG SUBSYSTEM: Media probing and transcoding ---
@mcp.tool(name="ffmpeg_exec")
async def ffmpeg_exec(args:str,timeout:int=300)->str:
    """执行FFmpeg命令(不含ffmpeg前缀)
    args: FFmpeg参数字符串
    例: "-i input.mp4 -ss 00:01:00 -t 30 -c copy clip.mp4"
    例: "-i video.mp4 -vf scale=1280:720 output.mp4"
    例: "-i input.mp4 -vn -acodec libmp3lame audio.mp3" """
    try:
        cmd=f'"{FFMPEG}" {args}'
        o,e,c=_run(cmd,timeout=timeout,shell=True)
        return J(stdout=o,stderr=e[-2000:],code=c)
    except Exception as e:
        return J(False,err=str(e))
@mcp.tool(name="ffmpeg_info")
async def ffmpeg_info(path:str)->str:
    """获取媒体文件信息(时长/分辨率/编码等)"""
    try:
        ffprobe=FFMPEG.replace("ffmpeg","ffprobe") if "ffmpeg" in FFMPEG.lower() else "ffprobe"
        import re as _re
        r=subprocess.run([ffprobe,"-v","quiet","-print_format","json","-show_format","-show_streams",str(R(path))],
                         capture_output=True,text=True,timeout=30,stdin=subprocess.DEVNULL)
        o=_re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\ufffd]','',r.stdout or "")
        info=json.loads(o) if o.strip() else {}
        fmt=info.get("format",{})
        streams=info.get("streams",[])
        result={"duration":fmt.get("duration"),"size":fmt.get("size"),
                "format":fmt.get("format_name"),"bitrate":fmt.get("bit_rate")}
        for s in streams:
            if s["codec_type"]=="video":
                result["video"]={"codec":s.get("codec_name"),"width":s.get("width"),
                                 "height":s.get("height"),"fps":s.get("r_frame_rate")}
            elif s["codec_type"]=="audio":
                result["audio"]={"codec":s.get("codec_name"),"sample_rate":s.get("sample_rate"),
                                 "channels":s.get("channels")}
        return J(**result)
    except Exception as e:
        return J(False,err=str(e))
@mcp.tool(name="ffmpeg_convert")
async def ffmpeg_convert(input:str,output:str,options:str="")->str:
    """媒体格式转换
    input: 输入文件  output: 输出文件
    options: 额外FFmpeg参数(如 "-crf 23 -preset fast")
    支持: mp4/avi/mkv/mov/mp3/wav/flac/gif/webm等互转"""
    try:
        cmd=f'"{FFMPEG}" -y -i "{R(input)}" {options} "{R(output)}"'
        o,e,c=_run(cmd,shell=True,timeout=300)
        return J(path=str(R(output)),code=c,log=e[-500:])
    except Exception as e:
        return J(False,err=str(e))
@mcp.tool(name="ffmpeg_clip")
async def ffmpeg_clip(input:str,output:str,start:str="00:00:00",
                      duration:str="",end:str="")->str:
    """视频/音频剪切
    start: 开始时间 "HH:MM:SS" 或秒数
    duration: 持续时长  end: 结束时间(二选一)"""
    try:
        cmd=f'"{FFMPEG}" -y -i "{R(input)}" -ss {start}'
        if duration:
            cmd+=f" -t {duration}"
        elif end:
            cmd+=f" -to {end}"
        cmd+=f' -c copy "{R(output)}"'
        o,e,c=_run(cmd,shell=True,timeout=120)
        return J(path=str(R(output)),code=c)
    except Exception as e:
        return J(False,err=str(e))
@mcp.tool(name="ffmpeg_screenshot")
async def ffmpeg_screenshot(input:str,output:str="screenshot.png",
                            time:str="00:00:01")->str:
    """从视频截取帧"""
    try:
        cmd=f'"{FFMPEG}" -y -ss {time} -i "{R(input)}" -vframes 1 "{R(output)}"'
        o,e,c=_run(cmd,shell=True,timeout=30)
        return J(path=str(R(output)),code=c)
    except Exception as e:
        return J(False,err=str(e))
@mcp.tool(name="ffmpeg_gif")
async def ffmpeg_gif(input:str,output:str="output.gif",fps:int=10,
                     width:int=480,start:str="00:00:00",duration:str="5")->str:
    """视频转GIF"""
    try:
        cmd=f'"{FFMPEG}" -y -ss {start} -t {duration} -i "{R(input)}" -vf "fps={fps},scale={width}:-1:flags=lanczos" "{R(output)}"'
        o,e,c=_run(cmd,shell=True,timeout=120)
        return J(path=str(R(output)),code=c)
    except Exception as e:
        return J(False,err=str(e))
# --- GIMP SUBSYSTEM: Batch image scripting interface ---
@mcp.tool(name="gimp_exec")
async def gimp_exec(script:str,timeout:int=120)->str:
    """执行GIMP Script-Fu脚本(批处理模式,通过stdin管道传递)
    script: Script-Fu代码
    例: '(let* ((img (car (gimp-file-load RUN-NONINTERACTIVE "in.png" "in.png")))
              (drawable (car (gimp-image-get-active-drawable img))))
          (plug-in-gauss RUN-NONINTERACTIVE img drawable 10 10 0)
          (file-png-save RUN-NONINTERACTIVE img drawable "out.png" "out.png" 0 9 1 1 1 1 1)
          (gimp-image-delete img))'"""
    try:
        cmd=[GIMP,"-i","--batch-interpreter","plug-in-script-fu-eval","-b","-"]
        full=f'{script}\n(gimp-quit 0)\n'
        r=subprocess.run(cmd,input=full,capture_output=True,text=True,
                         timeout=timeout,cwd=str(WD),creationflags=CF,errors="replace")
        o=r.stdout or ""
        e=r.stderr or ""
        if len(o)>4000:
            o=o[:2000]+"\n...(truncated)...\n"+o[-2000:]
        if len(e)>2000:
            e=e[-2000:]
        return J(stdout=o,stderr=e,code=r.returncode)
    except subprocess.TimeoutExpired:
        return J(False,err=f"GIMP超时({timeout}s)")
    except Exception as e:
        return J(False,err=str(e))
@mcp.tool(name="gimp_python")
async def gimp_python(script:str,timeout:int=120)->str:
    """执行GIMP Python-Fu脚本(批处理模式,通过stdin管道传递)
    script: Python代码(可用gimp/pdb模块)
    例: 加载图片、应用滤镜、导出"""
    try:
        full=f'(python-fu-eval RUN-NONINTERACTIVE 0 "{script.replace(chr(34),chr(92)+chr(34))}")\n(gimp-quit 0)\n'
        cmd=[GIMP,"-i","--batch-interpreter","plug-in-script-fu-eval","-b","-"]
        r=subprocess.run(cmd,input=full,capture_output=True,text=True,
                         timeout=timeout,cwd=str(WD),creationflags=CF,errors="replace")
        o=r.stdout or ""
        e=r.stderr or ""
        if len(o)>4000:
            o=o[:2000]+"\n...(truncated)...\n"+o[-2000:]
        if len(e)>2000:
            e=e[-2000:]
        return J(stdout=o,stderr=e,code=r.returncode)
    except subprocess.TimeoutExpired:
        return J(False,err=f"GIMP超时({timeout}s)")
    except Exception as e:
        return J(False,err=str(e))
# --- INKSCAPE SUBSYSTEM: Vector format conversion ---
@mcp.tool(name="inkscape_exec")
async def inkscape_exec(input:str="",actions:str="",output:str="",
                        export_type:str="")->str:
    """Inkscape命令行操作
    input: 输入SVG文件
    actions: Inkscape actions字符串(分号分隔)
      常用actions: select-all,object-align,export-filename,export-do
      例: "select-all;object-align:center page;export-filename:out.png;export-do"
    output: 导出路径  export_type: png|pdf|eps|svg|emf
    简单导出: inkscape_exec(input="a.svg",output="a.png",export_type="png")"""
    try:
        cmd=[INKSCAPE]
        if input:
            cmd.append(str(R(input)))
        if actions:
            cmd.extend(["--actions",actions])
        if output and export_type:
            cmd.extend([f"--export-type={export_type}",f"--export-filename={R(output)}"])
        elif output:
            cmd.extend([f"--export-filename={R(output)}"])
        o,e,c=_run(cmd,timeout=120)
        return J(stdout=o,stderr=e,code=c,output=str(R(output)) if output else "")
    except Exception as e:
        return J(False,err=str(e))
@mcp.tool(name="inkscape_convert")
async def inkscape_convert(input:str,output:str,dpi:int=300)->str:
    """Inkscape格式转换: SVG<->PDF/PNG/EPS/EMF
    高质量矢量转换,支持文字保留"""
    try:
        ext=Path(output).suffix.lstrip(".")
        cmd=[INKSCAPE,str(R(input)),f"--export-type={ext}",
             f"--export-filename={R(output)}",f"--export-dpi={dpi}"]
        o,e,c=_run(cmd,timeout=60)
        return J(path=str(R(output)),code=c)
    except Exception as e:
        return J(False,err=str(e))
# --- FREECAD SUBSYSTEM: Parametric CAD execution ---
@mcp.tool(name="freecad_exec")
async def freecad_exec(script:str,timeout:int=120)->str:
    """执行FreeCAD Python脚本(无GUI模式,使用FreeCADCmd.exe)
    script: FreeCAD Python代码(可用FreeCAD/Part/Mesh等模块)
    例: 创建零件、布尔运算、导出STL/STEP
    常用:
      import Part
      box = Part.makeBox(10,10,10)
      Part.show(box)
      doc = FreeCAD.ActiveDocument
      doc.saveCopy('output.FCStd')"""
    try:
        sf=WD/"_freecad.py"
        sf.write_text(script,encoding="utf-8")
        # FreeCADCmd.exe 直接接受Python脚本作为位置参数
        cmd=[FREECAD,str(sf)]
        o,e,c=_run(cmd,timeout=timeout)
        return J(stdout=o,stderr=e,code=c)
    except Exception as e:
        return J(False,err=str(e))
@mcp.tool(name="freecad_create")
async def freecad_create(shape:str="box",params:str="{}",
                         output:str="model.step",fmt:str="step")->str:
    """快捷创建FreeCAD零件并导出
    shape: box|cylinder|sphere|cone|torus
    params: {"length":10,"width":10,"height":10} (box)
            {"radius":5,"height":10} (cylinder/cone)
            {"radius":5} (sphere/torus, torus额外radius2)
    output: 输出文件  fmt: step|stl|obj|iges|brep"""
    kw=json.loads(params) if isinstance(params,str) else params
    op=str(R(output)).replace("\\","\\\\")
    lines=["import FreeCAD, Part"]
    lines.append("doc=FreeCAD.newDocument('Model')")
    if shape=="box":
        l,w,h=kw.get("length",10),kw.get("width",10),kw.get("height",10)
        lines.append(f"s=Part.makeBox({l},{w},{h})")
    elif shape=="cylinder":
        r,h=kw.get("radius",5),kw.get("height",10)
        lines.append(f"s=Part.makeCylinder({r},{h})")
    elif shape=="sphere":
        r=kw.get("radius",5)
        lines.append(f"s=Part.makeSphere({r})")
    elif shape=="cone":
        r1,r2,h=kw.get("radius",5),kw.get("radius2",0),kw.get("height",10)
        lines.append(f"s=Part.makeCone({r1},{r2},{h})")
    elif shape=="torus":
        r1,r2=kw.get("radius",10),kw.get("radius2",3)
        lines.append(f"s=Part.makeTorus({r1},{r2})")
    # 添加到文档并导出
    lines.append("obj=doc.addObject('Part::Feature','Shape')")
    lines.append("obj.Shape=s")
    lines.append("doc.recompute()")
    fmt_lower=fmt.lower()
    if fmt_lower=="step":
        lines.append(f"Part.export([obj],r'{op}')")
    elif fmt_lower=="stl":
        lines.append(f"import Mesh")
        lines.append(f"Mesh.export([obj],r'{op}')")
    elif fmt_lower=="iges":
        lines.append(f"Part.export([obj],r'{op}')")
    elif fmt_lower=="brep":
        lines.append(f"s.exportBrep(r'{op}')")
    else:
        lines.append(f"Part.export([obj],r'{op}')")
    lines.append(f"print('EXPORTED:',r'{op}')")
    return await freecad_exec("\n".join(lines))
# --- GODOT SUBSYSTEM: Project run and export ---
@mcp.tool(name="godot_exec")
async def godot_exec(project_path:str,script:str="",timeout:int=60)->str:
    """在Godot中执行操作
    project_path: Godot项目路径(含project.godot的目录)
    script: GDScript代码(写入临时文件执行)
      注意: 脚本需extends SceneTree, 用quit()退出
      例: 'extends SceneTree\nfunc _init():\n\tprint("hello")\n\tquit()'
    无script时返回项目信息"""
    try:
        pp=str(R(project_path))
        if script:
            # Godot 4脚本必须extends SceneTree才能headless运行
            if "extends" not in script:
                script=f"extends SceneTree\nfunc _init():\n\t"+script.replace("\n","\n\t")+"\n\tquit()"
            elif "quit()" not in script:
                script+="\n\tquit()"
            sf=WD/"_godot_script.gd"
            sf.write_text(script,encoding="utf-8")
            cmd=[GODOT,"--path",pp,"--headless","-s",str(sf)]
        else:
            cmd=[GODOT,"--path",pp,"--headless","--quit"]
        o,e,c=_run(cmd,timeout=timeout)
        return J(stdout=o,stderr=e,code=c)
    except Exception as e:
        return J(False,err=str(e))
@mcp.tool(name="godot_run")
async def godot_run(project_path:str,scene:str="",timeout:int=30)->str:
    """运行Godot项目(调试模式,自动超时退出)
    project_path: 项目路径  scene: 指定场景文件(.tscn)"""
    try:
        cmd=[GODOT,"--path",str(R(project_path))]
        if scene:
            cmd.append(scene)
        o,e,c=_run(cmd,timeout=timeout)
        return J(stdout=o,stderr=e,code=c)
    except subprocess.TimeoutExpired:
        return J(ok=True,msg=f"项目运行{timeout}秒后自动退出")
    except Exception as e:
        return J(False,err=str(e))
@mcp.tool(name="godot_export")
async def godot_export(project_path:str,preset:str="",output:str="")->str:
    """导出Godot项目
    preset: 导出预设名称(在export_presets.cfg中定义)
    output: 输出路径"""
    try:
        cmd=[GODOT,"--path",str(R(project_path)),"--headless",
             "--export-release",preset,str(R(output))]
        o,e,c=_run(cmd,timeout=300)
        return J(path=str(R(output)),code=c,stderr=e)
    except Exception as e:
        return J(False,err=str(e))
# --- UTILS SUBSYSTEM: Filesystem and command utilities ---
@mcp.tool(name="file_open")
async def file_open(path:str)->str:
    """用系统默认程序打开文件(Windows)"""
    try:
        p=str(R(path))
        os.startfile(p)
        return J(msg=f"已打开 {p}")
    except Exception as e:
        return J(False,err=str(e))
@mcp.tool(name="file_list")
async def file_list(dir:str="",pattern:str="",recursive:bool=False)->str:
    """列出目录文件
    pattern: 通配符如'*.py','*.pdf'
    recursive: 是否递归子目录"""
    try:
        d=R(dir) if dir else WD
        if pattern and recursive:
            files=[{"name":str(f.relative_to(d)),"size":f.stat().st_size,"dir":f.is_dir()}
                   for f in d.rglob(pattern)]
        elif pattern:
            files=[{"name":f.name,"size":f.stat().st_size,"dir":f.is_dir()}
                   for f in d.glob(pattern)]
        else:
            files=[{"name":f.name,"size":f.stat().st_size,"dir":f.is_dir()}
                   for f in sorted(d.iterdir())]
        return J(dir=str(d),files=files,count=len(files))
    except Exception as e:
        return J(False,err=str(e))
@mcp.tool(name="file_copy")
async def file_copy(src:str,dst:str)->str:
    """复制文件或目录"""
    try:
        s,d=str(R(src)),str(R(dst))
        if os.path.isdir(s):
            shutil.copytree(s,d)
        else:
            shutil.copy2(s,d)
        return J(src=s,dst=d)
    except Exception as e:
        return J(False,err=str(e))
@mcp.tool(name="file_move")
async def file_move(src:str,dst:str)->str:
    """移动/重命名文件"""
    try:
        s,d=str(R(src)),str(R(dst))
        shutil.move(s,d)
        return J(src=s,dst=d)
    except Exception as e:
        return J(False,err=str(e))
@mcp.tool(name="file_delete")
async def file_delete(path:str)->str:
    """删除文件或目录"""
    try:
        p=str(R(path))
        if os.path.isdir(p):
            shutil.rmtree(p)
        else:
            os.remove(p)
        return J(deleted=p)
    except Exception as e:
        return J(False,err=str(e))
@mcp.tool(name="file_read")
async def file_read(path:str,encoding:str="utf-8",lines:str="")->str:
    """读取文本文件
    lines: "1-10" 读取指定行范围, 空=全部"""
    try:
        p=str(R(path))
        text=Path(p).read_text(encoding=encoding)
        if lines:
            parts=lines.split("-")
            s,e=int(parts[0])-1,int(parts[-1])
            text="\n".join(text.splitlines()[s:e])
        if len(text)>10000:
            text=text[:5000]+"\n...(truncated)...\n"+text[-5000:]
        return J(path=p,content=text,length=len(text))
    except Exception as e:
        return J(False,err=str(e))
@mcp.tool(name="file_write")
async def file_write(path:str,content:str,encoding:str="utf-8",append:bool=False)->str:
    """写入文本文件"""
    try:
        p=str(R(path))
        mode="a" if append else "w"
        with open(p,mode,encoding=encoding) as f:
            f.write(content)
        return J(path=p,size=os.path.getsize(p))
    except Exception as e:
        return J(False,err=str(e))
@mcp.tool(name="run_cmd")
async def run_cmd(cmd:str,timeout:int=60,cwd:str="")->str:
    """执行系统命令
    cwd: 工作目录(可选,默认临时目录)"""
    try:
        o,e,c=_run(cmd,timeout=timeout,shell=True,cwd=str(R(cwd)) if cwd else None)
        return J(stdout=o,stderr=e,code=c)
    except Exception as e:
        return J(False,err=str(e))
@mcp.tool(name="run_python")
async def run_python(script:str,timeout:int=60)->str:
    """执行Python脚本(使用系统Python)"""
    try:
        sf=WD/"_script.py"
        sf.write_text(script,encoding="utf-8")
        import sys
        o,e,c=_run([sys.executable,str(sf)],timeout=timeout)
        return J(stdout=o,stderr=e,code=c)
    except Exception as e:
        return J(False,err=str(e))
@mcp.tool(name="system_info")
async def system_info()->str:
    """获取系统信息和已配置的软件路径"""
    import platform
    info={"os":platform.system(),"version":platform.version(),
          "arch":platform.architecture()[0],"python":platform.python_version(),
          "work_dir":str(WD),
          "tools":{
              "blender":BLENDER,"matlab":MATLAB,"ffmpeg":FFMPEG,
              "gimp":GIMP,"inkscape":INKSCAPE,"freecad":FREECAD,"godot":GODOT
          }}
    # 检测哪些可用
    for name,path in info["tools"].items():
        info["tools"][name]={"path":path,"found":os.path.isfile(path) if os.path.sep in path else shutil.which(path) is not None}
    return J(**info)
# --- ENTRY: MCP service bootstrap ---
# Keep stdio transport for compatibility with VS Code MCP clients.
if __name__ == "__main__":
    mcp.run(transport="stdio")
